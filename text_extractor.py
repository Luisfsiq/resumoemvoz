"""
Módulo para extração de texto de diferentes formatos de arquivo.
Suporta: .txt, .docx, .pptx, .pdf, .xlsx
"""

import os


def extract_from_txt(filepath: str) -> str:
    """Extrai texto de arquivos .txt"""
    encodings = ['utf-8', 'latin-1', 'cp1252']
    for encoding in encodings:
        try:
            with open(filepath, 'r', encoding=encoding) as f:
                return f.read()
        except (UnicodeDecodeError, UnicodeError):
            continue
    raise ValueError("Não foi possível decodificar o arquivo de texto.")


def extract_from_docx(filepath: str) -> str:
    """Extrai texto de arquivos .docx (Microsoft Word)"""
    from docx import Document
    doc = Document(filepath)
    paragraphs = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append(text)
    return '\n'.join(paragraphs)


def extract_from_pptx(filepath: str) -> str:
    """Extrai texto de arquivos .pptx (Microsoft PowerPoint)"""
    from pptx import Presentation
    prs = Presentation(filepath)
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        slide_content = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_content.append(text)
        if slide_content:
            slides_text.append(f"Slide {i}: " + '. '.join(slide_content))
    return '\n'.join(slides_text)


def extract_from_pdf(filepath: str) -> str:
    """Extrai texto de arquivos .pdf"""
    from PyPDF2 import PdfReader
    reader = PdfReader(filepath)
    pages_text = []
    for page in reader.pages:
        text = page.extract_text()
        if text and text.strip():
            pages_text.append(text.strip())
    return '\n'.join(pages_text)


def extract_from_xlsx(filepath: str) -> str:
    """Extrai texto de arquivos .xlsx (Microsoft Excel)"""
    from openpyxl import load_workbook
    wb = load_workbook(filepath, read_only=True, data_only=True)
    sheets_text = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows_text = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(cell) for cell in row if cell is not None]
            if cells:
                rows_text.append(' | '.join(cells))
        if rows_text:
            sheets_text.append(f"Planilha '{sheet_name}':\n" + '\n'.join(rows_text))
    wb.close()
    return '\n\n'.join(sheets_text)


# Mapeamento de extensões para funções de extração
EXTRACTORS = {
    '.txt': extract_from_txt,
    '.docx': extract_from_docx,
    '.pptx': extract_from_pptx,
    '.pdf': extract_from_pdf,
    '.xlsx': extract_from_xlsx,
}

ALLOWED_EXTENSIONS = set(EXTRACTORS.keys())


def get_supported_extensions() -> list:
    """Retorna lista de extensões suportadas"""
    return sorted(ALLOWED_EXTENSIONS)


def is_allowed_file(filename: str) -> bool:
    """Verifica se o arquivo tem uma extensão suportada"""
    ext = os.path.splitext(filename)[1].lower()
    return ext in ALLOWED_EXTENSIONS


def extract_text(filepath: str) -> str:
    """
    Extrai texto de um arquivo baseado na sua extensão.
    Retorna o texto extraído ou levanta uma exceção.
    """
    ext = os.path.splitext(filepath)[1].lower()
    if ext not in EXTRACTORS:
        raise ValueError(
            f"Formato '{ext}' não suportado. "
            f"Formatos aceitos: {', '.join(get_supported_extensions())}"
        )
    
    try:
        text = EXTRACTORS[ext](filepath)
        if not text or not text.strip():
            raise ValueError("Nenhum texto encontrado no arquivo.")
        return text.strip()
    except ValueError:
        raise
    except Exception as e:
        raise ValueError(f"Erro ao processar o arquivo: {str(e)}")
